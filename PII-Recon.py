import os
import threading
import time
import queue
import csv
import json
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# For Excel export, we use pandas if available.
try:
    import pandas as pd
except ImportError:
    pd = None

# Import the Presidio AnalyzerEngine
from presidio_analyzer import AnalyzerEngine

# Additional imports for PDFs, images, and OCR with easyocr
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from PIL import Image
except ImportError:
    Image = None

# Remove pytesseract dependency and use easyocr instead
try:
    import easyocr
except ImportError:
    easyocr = None

# Allowed file extensions for scanning (including common document types)
ALLOWED_EXTENSIONS = {
    ".txt", ".csv", ".json", ".log", ".xml",
    ".html", ".md", ".py", ".java", ".js", ".ini",
    ".pdf", ".png", ".jpg", ".jpeg", ".bmp", ".tiff",
    ".docx", ".xlsx", ".xls", ".pptx", ".ppt"
}

class PiiScannerApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PII Recon by Ru7")
        self.geometry("900x750")
        self.resizable(True, True)  # Window can now be resized and maximized
        
        # Initialize Presidio AnalyzerEngine (default built-in recognizers are loaded)
        self.engine = AnalyzerEngine()

        # Variables for scan control and results
        self.scan_thread = None
        self.pause_event = threading.Event()  # When set, scanning pauses
        self.stop_event = threading.Event()   # When set, scanning stops
        self.log_queue = queue.Queue()
        self.scan_results = []  # List to hold detection result dictionaries
        self.summary = {}       # Dict to hold counts per entity type
        self.total_files = 0
        self.scanned_files = 0

        # GUI state variables
        self.mode_var = tk.StringVar(value="file")
        self.threshold_var = tk.DoubleVar(value=0.5)
        self.selected_paths = []  # List of selected file paths (or folder)

        self.create_widgets()
        self.periodic_check()

    def create_widgets(self):
        # Frame for mode selection and file/folder choosing
        control_frame = ttk.Frame(self)
        control_frame.pack(pady=10, padx=10, fill=tk.X)

        ttk.Label(control_frame, text="Mode:").grid(row=0, column=0, sticky=tk.W, padx=5)
        modes = [("Single/Multiple File(s)", "file"),
                 ("Folder (recursive)", "folder")]
        col = 1
        for text, mode in modes:
            rb = ttk.Radiobutton(control_frame, text=text, variable=self.mode_var, value=mode)
            rb.grid(row=0, column=col, padx=5)
            col += 1

        self.select_button = ttk.Button(control_frame, text="Select File(s)/Folder", command=self.select_path)
        self.select_button.grid(row=0, column=col, padx=5)

        ttk.Label(control_frame, text="Min Confidence Threshold:").grid(row=1, column=0, sticky=tk.W, padx=5, pady=5)
        self.threshold_entry = ttk.Entry(control_frame, textvariable=self.threshold_var, width=10)
        self.threshold_entry.grid(row=1, column=1, sticky=tk.W, padx=5)

        self.start_button = ttk.Button(control_frame, text="Start Scan", command=self.start_scan)
        self.start_button.grid(row=1, column=2, padx=5)
        self.pause_button = ttk.Button(control_frame, text="Pause", command=self.pause_resume_scan, state=tk.DISABLED)
        self.pause_button.grid(row=1, column=3, padx=5)
        self.stop_button = ttk.Button(control_frame, text="Stop", command=self.stop_scan, state=tk.DISABLED)
        self.stop_button.grid(row=1, column=4, padx=5)

        self.progress_bar = ttk.Progressbar(self, orient="horizontal", length=850, mode="determinate")
        self.progress_bar.pack(pady=10)

        log_frame = ttk.LabelFrame(self, text="Live Scan Log")
        log_frame.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)
        self.log_text = tk.Text(log_frame, height=20, wrap=tk.NONE)
        self.log_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        log_scroll = ttk.Scrollbar(log_frame, orient="vertical", command=self.log_text.yview)
        log_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.configure(yscrollcommand=log_scroll.set)

        self.summary_label = ttk.Label(self, text="Summary: No scan yet")
        self.summary_label.pack(pady=5)

        # Action buttons: Only Export and Add Custom Recognizer
        action_frame = ttk.Frame(self)
        action_frame.pack(pady=10)
        self.export_button = ttk.Button(action_frame, text="Export Results", command=self.export_results, state=tk.DISABLED)
        self.export_button.grid(row=0, column=0, padx=5)
        self.custom_regex_button = ttk.Button(action_frame, text="Add Custom Recognizer", command=self.open_custom_regex_window)
        self.custom_regex_button.grid(row=0, column=1, padx=5)

    def select_path(self):
        mode = self.mode_var.get()
        if mode == "folder":
            folder = filedialog.askdirectory(title="Select Folder")
            if folder:
                self.selected_paths = [folder]
                self.log_queue.put(f"Selected folder: {folder}")
        else:
            files = filedialog.askopenfilenames(
                title="Select File(s)",
                filetypes=[("Supported Files", 
                            "*.txt *.csv *.json *.log *.xml *.html *.md *.py *.java *.js *.ini "
                            "*.pdf *.png *.jpg *.jpeg *.bmp *.tiff *.docx *.xlsx *.xls *.pptx *.ppt"),
                           ("All Files", "*.*")]
            )
            if files:
                self.selected_paths = list(files)
                self.log_queue.put(f"Selected {len(files)} file(s)")

    def start_scan(self):
        if not self.selected_paths:
            messagebox.showwarning("No Selection", "Please select file(s) or folder to scan.")
            return

        try:
            threshold = float(self.threshold_var.get())
        except ValueError:
            messagebox.showerror("Invalid Threshold", "Please enter a valid number for confidence threshold.")
            return

        # Build list of files to scan based on allowed extensions
        file_list = []
        if self.mode_var.get() == "folder":
            folder = self.selected_paths[0]
            for root, dirs, files in os.walk(folder):
                for file in files:
                    if os.path.splitext(file)[1].lower() in ALLOWED_EXTENSIONS:
                        file_list.append(os.path.join(root, file))
        else:
            for f in self.selected_paths:
                if os.path.splitext(f)[1].lower() in ALLOWED_EXTENSIONS:
                    file_list.append(f)
                else:
                    self.log_queue.put(f"Skipping unsupported file type: {f}")

        if not file_list:
            messagebox.showinfo("No Valid Files", "No valid files found for scanning.")
            return

        self.total_files = len(file_list)
        self.scanned_files = 0
        self.scan_results = []
        self.summary = {}
        self.progress_bar["value"] = 0
        self.progress_bar["maximum"] = self.total_files

        # Disable buttons during scan
        self.start_button.config(state=tk.DISABLED)
        self.select_button.config(state=tk.DISABLED)
        self.export_button.config(state=tk.DISABLED)
        self.custom_regex_button.config(state=tk.DISABLED)
        self.pause_button.config(state=tk.NORMAL)
        self.stop_button.config(state=tk.NORMAL)

        self.log_text.delete(1.0, tk.END)
        self.log_queue.put("Starting scan...\n")
        
        self.pause_event.clear()
        self.stop_event.clear()

        # Start scanning in a separate thread
        self.scan_thread = threading.Thread(target=self.scan_worker, args=(file_list, threshold), daemon=True)
        self.scan_thread.start()

    def scan_worker(self, file_list, threshold):
        for file_path in file_list:
            if self.stop_event.is_set():
                self.log_queue.put("Scan stopped by user.\n")
                break

            self.log_queue.put(f"Scanning file: {file_path}")
            lines = self.extract_text(file_path)
            if not lines:
                self.log_queue.put(f"No text extracted from {file_path}.\n")
                self.scanned_files += 1
                self.progress_bar["value"] = self.scanned_files
                continue

            for i, line in enumerate(lines, start=1):
                while self.pause_event.is_set() and not self.stop_event.is_set():
                    time.sleep(0.1)
                if self.stop_event.is_set():
                    break

                results = self.engine.analyze(text=line, language="en", score_threshold=threshold)
                for res in results:
                    record = {
                        "file_path": file_path,
                        "file_name": os.path.basename(file_path),
                        "line_number": i,
                        "entity_type": res.entity_type,
                        "matched_text": line[res.start:res.end],
                        "confidence_score": res.score
                    }
                    self.scan_results.append(record)
                    self.summary[res.entity_type] = self.summary.get(res.entity_type, 0) + 1
                    self.log_queue.put(
                        f"File: {os.path.basename(file_path)} | Line {i} | Entity: {res.entity_type} | "
                        f"Text: {line[res.start:res.end].strip()} | Score: {res.score:.2f}"
                    )

            self.scanned_files += 1
            self.progress_bar["value"] = self.scanned_files
            self.update_summary()

        self.log_queue.put("Scan completed.\n")
        self.export_button.config(state=tk.NORMAL)
        self.start_button.config(state=tk.NORMAL)
        self.select_button.config(state=tk.NORMAL)
        self.custom_regex_button.config(state=tk.NORMAL)
        self.pause_button.config(state=tk.DISABLED)
        self.stop_button.config(state=tk.DISABLED)

    def extract_text(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.extract_text_from_pdf(file_path)
        elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}:
            return self.extract_text_from_image(file_path)
        elif ext == ".docx":
            return self.extract_text_from_docx(file_path)
        elif ext == ".xlsx":
            return self.extract_text_from_excel(file_path)
        elif ext == ".xls":
            return self.extract_text_from_excel_xls(file_path)
        elif ext == ".pptx":
            return self.extract_text_from_pptx(file_path)
        elif ext == ".ppt":
            self.log_queue.put("PPT format is not supported. Please convert to PPTX.")
            return []
        else:
            try:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    return f.readlines()
            except Exception as e:
                self.log_queue.put(f"Error reading file {file_path}: {e}")
                return []

    def extract_text_from_pdf(self, file_path):
        if PyPDF2 is None:
            self.log_queue.put("PyPDF2 is not installed. Cannot process PDF files.")
            return []
        text = ""
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.splitlines()
        except Exception as e:
            self.log_queue.put(f"Error processing PDF file {file_path}: {e}")
            return []

    def extract_text_from_image(self, file_path):
        """
        Extract text from an image file using easyocr.
        """
        if easyocr is None:
            self.log_queue.put("easyocr is not installed. Cannot process image files.")
            return []
        try:
            # Initialize the easyocr reader (using English; set gpu=True if you have GPU support)
            reader = easyocr.Reader(['en'], gpu=False)
            result = reader.readtext(file_path, detail=0)
            # result is a list of strings extracted from the image
            return result
        except Exception as e:
            self.log_queue.put(f"Error processing image file {file_path}: {e}")
            return []

    def extract_text_from_docx(self, file_path):
        try:
            import docx
        except ImportError:
            self.log_queue.put("python-docx is not installed. Cannot process DOCX files.")
            return []
        try:
            doc = docx.Document(file_path)
            lines = [para.text for para in doc.paragraphs if para.text.strip()]
            return lines
        except Exception as e:
            self.log_queue.put(f"Error processing DOCX file {file_path}: {e}")
            return []

    def extract_text_from_excel(self, file_path):
        try:
            import openpyxl
        except ImportError:
            self.log_queue.put("openpyxl is not installed. Cannot process XLSX files.")
            return []
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            lines = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        lines.append(row_text)
            return lines
        except Exception as e:
            self.log_queue.put(f"Error processing Excel file {file_path}: {e}")
            return []

    def extract_text_from_excel_xls(self, file_path):
        try:
            import xlrd
        except ImportError:
            self.log_queue.put("xlrd is not installed. Cannot process .xls files.")
            return []
        try:
            wb = xlrd.open_workbook(file_path)
            lines = []
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row(row_idx)
                    row_text = " ".join([str(cell.value) for cell in row])
                    if row_text.strip():
                        lines.append(row_text)
            return lines
        except Exception as e:
            self.log_queue.put(f"Error processing Excel file {file_path}: {e}")
            return []

    def extract_text_from_pptx(self, file_path):
        try:
            from pptx import Presentation
        except ImportError:
            self.log_queue.put("python-pptx is not installed. Cannot process PPTX files.")
            return []
        try:
            prs = Presentation(file_path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.extend(shape.text.splitlines())
            return lines
        except Exception as e:
            self.log_queue.put(f"Error processing PPTX file {file_path}: {e}")
            return []

    def update_summary(self):
        summary_text = "Summary: " + ", ".join([f"{k}: {v}" for k, v in self.summary.items()]) if self.summary else "No PII found yet."
        self.summary_label.config(text=summary_text)

    def pause_resume_scan(self):
        if not self.pause_event.is_set():
            self.pause_event.set()
            self.pause_button.config(text="Resume")
            self.log_queue.put("Scan paused.")
        else:
            self.pause_event.clear()
            self.pause_button.config(text="Pause")
            self.log_queue.put("Scan resumed.")

    def stop_scan(self):
        if messagebox.askyesno("Stop Scan", "Are you sure you want to stop the scan?"):
            self.stop_event.set()
            self.log_queue.put("Stopping scan...")

    def export_results(self):
        if not self.scan_results:
            messagebox.showinfo("No Results", "There are no scan results to export.")
            return

        filetypes = [("CSV Files", "*.csv")]
        if pd is not None:
            filetypes.append(("Excel Files", "*.xlsx"))
        export_path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=filetypes)
        if not export_path:
            return

        try:
            if export_path.endswith(".xlsx") and pd is not None:
                df = pd.DataFrame(self.scan_results)
                df.to_excel(export_path, index=False)
            else:
                with open(export_path, mode="w", newline="", encoding="utf-8") as f:
                    writer = csv.DictWriter(f, fieldnames=["file_path", "file_name", "line_number", "entity_type", "matched_text", "confidence_score"])
                    writer.writeheader()
                    writer.writerows(self.scan_results)
            messagebox.showinfo("Export Successful", f"Results exported to {export_path}")
        except Exception as e:
            messagebox.showerror("Export Error", f"Failed to export results: {e}")

    def open_custom_regex_window(self):
        """Open a window to add multiple custom regex recognizers."""
        win = tk.Toplevel(self)
        win.title("Add Custom Recognizers")
        win.geometry("500x300")
        
        # Frame for input fields
        input_frame = ttk.Frame(win)
        input_frame.pack(pady=10, padx=10, fill=tk.X)
        
        ttk.Label(input_frame, text="Entity Type:").grid(row=0, column=0, sticky=tk.W, pady=5)
        entity_entry = ttk.Entry(input_frame, width=40)
        entity_entry.grid(row=0, column=1, pady=5)
        
        ttk.Label(input_frame, text="Regex Pattern:").grid(row=1, column=0, sticky=tk.W, pady=5)
        regex_entry = ttk.Entry(input_frame, width=40)
        regex_entry.grid(row=1, column=1, pady=5)
        
        ttk.Label(input_frame, text="Confidence Score (e.g., 0.5):").grid(row=2, column=0, sticky=tk.W, pady=5)
        score_entry = ttk.Entry(input_frame, width=10)
        score_entry.insert(0, "0.5")
        score_entry.grid(row=2, column=1, sticky=tk.W, pady=5)
        
        # Listbox to show added custom recognizers
        listbox_frame = ttk.Frame(win)
        listbox_frame.pack(pady=10, padx=10, fill=tk.BOTH, expand=True)
        ttk.Label(listbox_frame, text="Custom Recognizers Added:").pack(anchor=tk.W)
        recognizer_list = tk.Listbox(listbox_frame, height=5)
        recognizer_list.pack(fill=tk.BOTH, expand=True)
        
        def add_recognizer():
            entity = entity_entry.get().strip()
            regex = regex_entry.get().strip()
            try:
                score = float(score_entry.get().strip())
            except ValueError:
                messagebox.showerror("Invalid Score", "Please enter a valid number for confidence score.")
                return
            if not entity or not regex:
                messagebox.showerror("Missing Data", "Entity type and regex pattern cannot be empty.")
                return
            # Add the custom recognizer and update the listbox
            self.add_custom_recognizer(entity, regex, score)
            recognizer_list.insert(tk.END, f"{entity.upper()} | {regex} | {score}")
            # Clear fields for next input
            entity_entry.delete(0, tk.END)
            regex_entry.delete(0, tk.END)
            score_entry.delete(0, tk.END)
            score_entry.insert(0, "0.5")
        
        add_button = ttk.Button(win, text="Add Recognizer", command=add_recognizer)
        add_button.pack(pady=5)
        
        done_button = ttk.Button(win, text="Done", command=win.destroy)
        done_button.pack(pady=5)

    def add_custom_recognizer(self, entity, regex, score):
        """
        Add a custom recognizer to the Presidio AnalyzerEngine using the provided entity, regex, and score.
        """
        try:
            from presidio_analyzer import Pattern, PatternRecognizer
            pattern = Pattern(name=entity, regex=regex, score=score)
            recognizer = PatternRecognizer(supported_entity=entity.upper(), patterns=[pattern])
            self.engine.registry.add_recognizer(recognizer)
            self.log_queue.put(f"Custom recognizer for '{entity}' added.")
        except Exception as e:
            self.log_queue.put(f"Error adding custom recognizer: {e}")

    def periodic_check(self):
        """Periodically update the log text widget with messages from the queue."""
        try:
            while True:
                msg = self.log_queue.get_nowait()
                self.log_text.insert(tk.END, msg + "\n")
                self.log_text.see(tk.END)
        except queue.Empty:
            pass
        self.after(100, self.periodic_check)

    def on_closing(self):
        if self.scan_thread and self.scan_thread.is_alive():
            if messagebox.askyesno("Quit", "A scan is in progress. Do you really want to exit?"):
                self.stop_event.set()
                self.destroy()
        else:
            self.destroy()

if __name__ == '__main__':
    app = PiiScannerApp()
    app.protocol("WM_DELETE_WINDOW", app.on_closing)
    app.mainloop()
